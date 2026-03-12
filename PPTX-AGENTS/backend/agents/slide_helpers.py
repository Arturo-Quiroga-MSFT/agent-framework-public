"""
agents/slide_helpers.py
-----------------------
Design-aware slide-builder helpers injected into the generated code's exec
namespace.  The code model calls these functions instead of raw python-pptx,
so every slide gets consistent, polished formatting without the model needing
to know any python-pptx internals.

Available in generated code (no imports needed):
    make_presentation()
    add_title_slide(prs, title, subtitle, notes, theme_name)
    add_content_slide(prs, title, bullets, notes, theme_name)
    add_two_column_slide(prs, title, left_title, left_bullets,
                         right_title, right_bullets, notes, theme_name)
    add_section_break_slide(prs, title, subtitle, notes, theme_name)
    add_references_slide(prs, title, links, notes, theme_name)
    add_closing_slide(prs, title, cta, notes, theme_name)

THEME and OUTPUT_PATH are pre-set variables — no need to hardcode them.
"""

from __future__ import annotations

import io
from urllib.parse import urlparse

import requests as _requests
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Allowed image hosts (security gate) ─────────────────────────────────────
ALLOWED_IMAGE_HOSTS: frozenset[str] = frozenset({
    "learn.microsoft.com",
    "docs.microsoft.com",
    "azure.microsoft.com",
    "devblogs.microsoft.com",
    "techcommunity.microsoft.com",
    "cdn-dynmedia-1.microsoft.com",
    "msazure.blob.core.windows.net",
})


def _fetch_image_bytes(url: str) -> io.BytesIO | None:
    """Download an image URL, enforcing the allowed-host list. Returns BytesIO or None."""
    try:
        host = urlparse(url).hostname or ""
        # Accept exact match or subdomain of allowed hosts
        if not any(host == h or host.endswith("." + h) for h in ALLOWED_IMAGE_HOSTS):
            return None
        r = _requests.get(url, timeout=8, headers={"User-Agent": "PPTX-Agents/1.0"})
        if r.status_code == 200 and r.content:
            return io.BytesIO(r.content)
    except Exception:
        pass
    return None

# ── Slide canvas (widescreen 16:9) ─────────────────────────────────────────
_W = Inches(13.33)
_H = Inches(7.5)

# ── Theme palette definitions ───────────────────────────────────────────────
THEMES: dict[str, dict] = {
    "professional": dict(
        bg         = RGBColor(0x00, 0x33, 0x66),   # navy
        bg_content = RGBColor(0xFF, 0xFF, 0xFF),   # white
        accent     = RGBColor(0x00, 0x7A, 0xC3),   # azure blue
        title_txt  = RGBColor(0xFF, 0xFF, 0xFF),
        body_txt   = RGBColor(0x1A, 0x1A, 0x2E),
        sub_txt    = RGBColor(0xB3, 0xD4, 0xF0),   # pale blue
        col_hdr    = RGBColor(0x00, 0x33, 0x66),
    ),
    "modern": dict(
        bg         = RGBColor(0x1A, 0x1A, 0x2E),
        bg_content = RGBColor(0xF4, 0xF5, 0xF7),
        accent     = RGBColor(0x6C, 0x63, 0xFF),
        title_txt  = RGBColor(0xFF, 0xFF, 0xFF),
        body_txt   = RGBColor(0x1A, 0x1A, 0x2E),
        sub_txt    = RGBColor(0xC0, 0xBB, 0xFF),
        col_hdr    = RGBColor(0x1A, 0x1A, 0x2E),
    ),
    "minimal": dict(
        bg         = RGBColor(0xFF, 0xFF, 0xFF),
        bg_content = RGBColor(0xFF, 0xFF, 0xFF),
        accent     = RGBColor(0x00, 0x00, 0x00),
        title_txt  = RGBColor(0x00, 0x00, 0x00),
        body_txt   = RGBColor(0x33, 0x33, 0x33),
        sub_txt    = RGBColor(0x66, 0x66, 0x66),
        col_hdr    = RGBColor(0x22, 0x22, 0x22),
    ),
    "bold": dict(
        bg         = RGBColor(0x4A, 0x00, 0x80),   # deep purple
        bg_content = RGBColor(0xFF, 0xFF, 0xFF),
        accent     = RGBColor(0xFF, 0x6B, 0x00),   # vivid orange
        title_txt  = RGBColor(0xFF, 0xFF, 0xFF),
        body_txt   = RGBColor(0x1A, 0x1A, 0x2E),
        sub_txt    = RGBColor(0xFF, 0xD0, 0x99),
        col_hdr    = RGBColor(0x4A, 0x00, 0x80),
    ),
}
_FALLBACK = "professional"


# ── Private utilities ───────────────────────────────────────────────────────

def _t(name: str) -> dict:
    return THEMES.get(name, THEMES[_FALLBACK])


def _fill_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _rect(slide, left: float, top: float, width: float, height: float,
          color: RGBColor):
    """Add a filled rectangle with no border (all dims in inches)."""
    shape = slide.shapes.add_shape(
        1,  # MSO_AUTO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _tf(slide, left: float, top: float, width: float, height: float):
    """Add a word-wrapped textbox; return its TextFrame."""
    tb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    tb.text_frame.word_wrap = True
    return tb.text_frame


def _p0(tf, text: str, size: int, bold: bool, color: RGBColor,
        align=PP_ALIGN.LEFT) -> None:
    """Set the first (pre-existing) paragraph of a fresh TextFrame."""
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    p.space_after = Pt(0)


def _pn(tf, text: str, size: int, bold: bool, color: RGBColor,
        align=PP_ALIGN.LEFT, space_before: int = 8) -> None:
    """Append a paragraph to an existing TextFrame."""
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    p.space_before = Pt(space_before)
    p.space_after = Pt(0)


def _notes(slide, text: str | None) -> None:
    if text:
        slide.notes_slide.notes_text_frame.text = str(text)


def _bullet_list(tf, bullets: list[str], size: int, color: RGBColor) -> None:
    """Render a bullet list into a TextFrame (first bullet uses _p0)."""
    if not bullets:
        return
    _p0(tf, f"\u2022  {bullets[0]}", size, False, color)
    for b in bullets[1:]:
        _pn(tf, f"\u2022  {b}", size, False, color, space_before=8)


# ── Public slide constructors ───────────────────────────────────────────────

def make_presentation() -> Presentation:
    """Return a blank widescreen (16:9) Presentation."""
    prs = Presentation()
    prs.slide_width  = _W
    prs.slide_height = _H
    return prs


def add_title_slide(prs: Presentation, title: str, subtitle: str = "",
                    notes: str = "", theme_name: str = "professional"):
    """
    Full-bleed dark opening slide:
        ┌─┬──────────────────────────────┐
        │ │                              │
        │█│   BIG TITLE                 │
        │ │   subtitle (lighter)         │
        │ │______________________________|
        └─┴──────────────────────────────┘
    """
    t = _t(theme_name)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    _fill_bg(slide, t["bg"])

    # Left accent bar
    _rect(slide, 0, 0, 0.45, 7.5, t["accent"])
    # Bottom horizontal rule
    _rect(slide, 0.45, 6.75, 12.88, 0.05, t["accent"])

    # Title
    tf = _tf(slide, 1.0, 2.0, 11.8, 2.2)
    _p0(tf, title, 42, True, t["title_txt"])

    # Subtitle
    if subtitle:
        tf2 = _tf(slide, 1.0, 4.4, 11.8, 1.5)
        _p0(tf2, subtitle, 22, False, t["sub_txt"])

    _notes(slide, notes)
    return slide


def add_content_slide(prs: Presentation, title: str, bullets: list[str] | None = None,
                      notes: str = "", theme_name: str = "professional"):
    """
    Light-bg slide with dark header band:
        ┌──────────────────────────────────┐
        │  TITLE (dark band, white text)   │
        │▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬   │
        │  • bullet                        │
        │  • bullet                        │
        └──────────────────────────────────┘
    """
    t = _t(theme_name)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    _fill_bg(slide, t["bg_content"])

    # Dark header band
    _rect(slide, 0, 0, 13.33, 1.28, t["bg"])
    # Accent underline
    _rect(slide, 0, 1.28, 13.33, 0.07, t["accent"])

    # Title
    tf = _tf(slide, 0.45, 0.13, 12.4, 1.0)
    _p0(tf, title, 30, True, t["title_txt"])

    # Bullet list
    if bullets:
        bft = _tf(slide, 0.6, 1.55, 12.1, 5.65)
        _bullet_list(bft, bullets, 20, t["body_txt"])

    _notes(slide, notes)
    return slide


def add_two_column_slide(prs: Presentation, title: str,
                         left_title: str = "", left_bullets: list[str] | None = None,
                         right_title: str = "", right_bullets: list[str] | None = None,
                         notes: str = "", theme_name: str = "professional"):
    """
    Two-column comparison slide:
        ┌──────────────────────────────────┐
        │  TITLE (dark header band)        │
        │▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬  │
        │  Left title   │  Right title     │
        │  • bullet     │  • bullet        │
        └──────────────────────────────────┘
    """
    t = _t(theme_name)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    _fill_bg(slide, t["bg_content"])
    _rect(slide, 0, 0, 13.33, 1.28, t["bg"])
    _rect(slide, 0, 1.28, 13.33, 0.07, t["accent"])

    # Slide title
    tf = _tf(slide, 0.45, 0.13, 12.4, 1.0)
    _p0(tf, title, 30, True, t["title_txt"])

    # Left column header
    if left_title:
        ltf = _tf(slide, 0.45, 1.52, 5.95, 0.6)
        _p0(ltf, left_title, 18, True, t["col_hdr"])
        _rect(slide, 0.45, 2.1, 5.95, 0.05, t["accent"])

    # Left bullets
    if left_bullets:
        lbf = _tf(slide, 0.55, 2.25, 5.85, 5.0)
        _bullet_list(lbf, left_bullets, 18, t["body_txt"])

    # Vertical divider
    _rect(slide, 6.64, 1.35, 0.05, 5.9, t["accent"])

    # Right column header
    if right_title:
        rtf = _tf(slide, 6.85, 1.52, 6.0, 0.6)
        _p0(rtf, right_title, 18, True, t["col_hdr"])
        _rect(slide, 6.85, 2.1, 6.0, 0.05, t["accent"])

    # Right bullets
    if right_bullets:
        rbf = _tf(slide, 6.95, 2.25, 5.9, 5.0)
        _bullet_list(rbf, right_bullets, 18, t["body_txt"])

    _notes(slide, notes)
    return slide


def add_section_break_slide(prs: Presentation, title: str, subtitle: str = "",
                             notes: str = "", theme_name: str = "professional"):
    """
    Full-bleed dark transition / section-break slide with wide left accent bar.
        ┌───┬─────────────────────────────┐
        │   │                             │
        │███│   SECTION TITLE             │
        │   │   subtitle                  │
        │   │─────────────────────────────│
        └───┴─────────────────────────────┘
    """
    t = _t(theme_name)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    _fill_bg(slide, t["bg"])

    # Wide left accent bar
    _rect(slide, 0, 0, 0.9, 7.5, t["accent"])
    # Horizontal underline below title area
    _rect(slide, 0.9, 5.3, 12.43, 0.05, t["sub_txt"])

    # Title
    tf = _tf(slide, 1.5, 2.1, 11.3, 2.2)
    _p0(tf, title, 38, True, t["title_txt"])

    # Subtitle
    if subtitle:
        tf2 = _tf(slide, 1.5, 4.5, 11.3, 1.4)
        _p0(tf2, subtitle, 20, False, t["sub_txt"])

    _notes(slide, notes)
    return slide


def add_content_slide_with_image(
        prs: Presentation, title: str,
        bullets: list[str] | None = None,
        image_url: str = "",
        image_caption: str = "",
        notes: str = "",
        theme_name: str = "professional"):
    """
    Content slide with an image panel on the right:
        ┌──────────────────────────────────┐
        │  TITLE (dark header band)        │
        │▇▇▇▇▇▇▇▇▇▇▇▇▇▇▇▇▇▇▇▇▇▇▇▇▇▇▇▇▇▇▇  │
        │  • bullet     │ [  image   ] │
        │  • bullet     │  caption     │
        └──────────────────────────────────┘
    image_url must be from ALLOWED_IMAGE_HOSTS; silently omitted otherwise.
    """
    t = _t(theme_name)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    _fill_bg(slide, t["bg_content"])
    _rect(slide, 0, 0, 13.33, 1.28, t["bg"])
    _rect(slide, 0, 1.28, 13.33, 0.07, t["accent"])

    # Title
    tf = _tf(slide, 0.45, 0.13, 12.4, 1.0)
    _p0(tf, title, 30, True, t["title_txt"])

    # Bullets on left portion
    if bullets:
        bft = _tf(slide, 0.6, 1.55, 6.5, 5.65)
        _bullet_list(bft, bullets, 18, t["body_txt"])

    # Image on right portion
    if image_url:
        img_bytes = _fetch_image_bytes(image_url)
        if img_bytes:
            try:
                slide.shapes.add_picture(
                    img_bytes,
                    Inches(7.3), Inches(1.55), Inches(5.6), Inches(4.9),
                )
            except Exception:
                img_bytes = None  # fall through to caption-only

            # Caption / source link
            if image_caption:
                cf = _tf(slide, 7.3, 6.5, 5.6, 0.6)
                _p0(cf, f"\u2192 {image_caption}", 9, False, t["accent"], PP_ALIGN.RIGHT)

    _notes(slide, notes)
    return slide


def add_references_slide(prs: Presentation, title: str = "Sources & References",
                         links: list | None = None,
                         notes: str = "", theme_name: str = "professional"):
    """
    References slide — compact list of clickable source URLs.
    links: list of {"text": "...", "url": "..."} dicts, or plain URL strings.
    """
    if links is None:
        links = []

    t = _t(theme_name)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_bg(slide, t["bg_content"])

    # Header band + accent underline
    _rect(slide, 0, 0, 13.33, 1.6, t["bg"])
    _rect(slide, 0, 1.6, 13.33, 0.06, t["accent"])

    # Title in header band
    tf = _tf(slide, 0.4, 0.2, 12.5, 1.2)
    _p0(tf, title, 28, True, t["title_txt"])

    # Link list
    if links:
        lf = _tf(slide, 0.7, 2.0, 12.0, 5.1)
        lf.word_wrap = True
        first = True
        for item in links:
            if isinstance(item, dict):
                display = item.get("text") or item.get("url", "")
                url     = item.get("url", "")
            else:
                display = url = str(item)

            if not url:
                continue

            p = lf.paragraphs[0] if first else lf.add_paragraph()
            first = False
            p.space_before = Pt(0 if p == lf.paragraphs[0] else 10)
            p.space_after  = Pt(0)

            # Arrow bullet
            bullet = p.add_run()
            bullet.text = "→  "
            bullet.font.size  = Pt(14)
            bullet.font.bold  = True
            bullet.font.color.rgb = t["accent"]

            # Display text — hyperlinked if URL present
            run = p.add_run()
            run.text = display
            run.font.size  = Pt(14)
            run.font.bold  = False
            run.font.color.rgb = t["accent"]
            if url.startswith("http"):
                run.hyperlink.address = url

            # URL in smaller dimmed text underneath (if display ≠ url)
            if display != url and url.startswith("http"):
                pu = lf.add_paragraph()
                pu.space_before = Pt(2)
                pu.space_after  = Pt(0)
                ru = pu.add_run()
                ru.text = f"    {url}"
                ru.font.size  = Pt(10)
                ru.font.bold  = False
                ru.font.color.rgb = RGBColor(0x80, 0x80, 0x80)
                ru.hyperlink.address = url

    _notes(slide, notes)
    return slide


def add_closing_slide(prs: Presentation, title: str, cta: str = "",
                      notes: str = "", theme_name: str = "professional"):
    """Closing / thank-you slide — same design as the opening title slide."""
    return add_title_slide(prs, title, cta, notes, theme_name)
