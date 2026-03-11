"""
shared/pptx_utils.py
--------------------
Common PPTX utilities shared by analyst and builder agents.

  - extract_text(path)     → markdown string via markitdown
  - render_slides(path)    → list of base64 PNG strings (one per slide)
  - slide_count(path)      → int
"""

from __future__ import annotations

import base64
import io
import os
from pathlib import Path

from pptx import Presentation


# ---------------------------------------------------------------------------
# Text extraction (markitdown)
# ---------------------------------------------------------------------------

def extract_text(pptx_path: str | Path) -> str:
    """Extract all text from a PPTX as a markdown string using markitdown."""
    try:
        from markitdown import MarkItDown
        md = MarkItDown()
        result = md.convert(str(pptx_path))
        return result.text_content
    except ImportError:
        # Fallback: plain text extraction via python-pptx
        return _extract_text_fallback(pptx_path)


def _extract_text_fallback(pptx_path: str | Path) -> str:
    """Plain-text extraction fallback when markitdown is not available."""
    prs = Presentation(str(pptx_path))
    lines: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"\n## Slide {i}\n")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        lines.append(text)
    return "\n".join(lines)


# ---------------------------------------------------------------------------
# Slide rendering → base64 PNG  (requires LibreOffice or pdf2image)
# ---------------------------------------------------------------------------

def render_slides_to_images(pptx_path: str | Path) -> list[str]:
    """
    Render each slide to a base64-encoded PNG string.

    Strategy (in order of preference):
      1. python-pptx + Pillow thumbnail (fast, shapes only — no fonts/images)
      2. LibreOffice headless → PDF → pdf2image (highest fidelity, requires LO)

    Returns a list of base64 PNG strings, one per slide.
    """
    return _render_with_libreoffice(pptx_path) or _render_with_pptx_thumbnail(pptx_path)


def _render_with_libreoffice(pptx_path: str | Path) -> list[str] | None:
    """Convert via LibreOffice headless → PDF → images. Returns None if unavailable."""
    import shutil, subprocess, tempfile
    if not shutil.which("libreoffice") and not shutil.which("soffice"):
        return None
    try:
        from pdf2image import convert_from_path
    except ImportError:
        return None

    lo_bin = shutil.which("libreoffice") or shutil.which("soffice")
    with tempfile.TemporaryDirectory() as tmpdir:
        subprocess.run(
            [lo_bin, "--headless", "--convert-to", "pdf", "--outdir", tmpdir, str(pptx_path)],
            check=True, capture_output=True,
        )
        pdf_files = list(Path(tmpdir).glob("*.pdf"))
        if not pdf_files:
            return None
        images = convert_from_path(str(pdf_files[0]), dpi=150)
        result = []
        for img in images:
            buf = io.BytesIO()
            img.save(buf, format="PNG")
            result.append(base64.b64encode(buf.getvalue()).decode())
        return result


def _render_with_pptx_thumbnail(pptx_path: str | Path) -> list[str]:
    """Lightweight fallback: render slide background colours + text as PIL images."""
    from PIL import Image, ImageDraw, ImageFont

    prs = Presentation(str(pptx_path))
    results: list[str] = []

    for slide in prs.slides:
        width, height = int(prs.slide_width.pt), int(prs.slide_height.pt)
        img = Image.new("RGB", (width, height), color=(255, 255, 255))
        draw = ImageDraw.Draw(img)

        y_offset = 20
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        draw.text((20, y_offset), text, fill=(30, 30, 30))
                        y_offset += 18

        buf = io.BytesIO()
        img.save(buf, format="PNG")
        results.append(base64.b64encode(buf.getvalue()).decode())

    return results


# ---------------------------------------------------------------------------
# Metadata helpers
# ---------------------------------------------------------------------------

def slide_count(pptx_path: str | Path) -> int:
    return len(Presentation(str(pptx_path)).slides)


def slide_titles(pptx_path: str | Path) -> list[str]:
    prs = Presentation(str(pptx_path))
    titles = []
    for slide in prs.slides:
        title_shape = slide.shapes.title
        titles.append(title_shape.text.strip() if title_shape and title_shape.has_text_frame else "")
    return titles
